import os
import re
import datetime
import logging

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

logger = logging.getLogger("breaknote")

def parse_date(date_str):
    try:
        return datetime.datetime.strptime(date_str, "%Y-%m-%d").date()
    except Exception:
        # Fallback to today if parsing fails
        return datetime.date.today()

def extract_text_pdf(path, limit=10):
    if not fitz:
        return ""
    text = ""
    try:
        doc = fitz.open(path)
        for i in range(min(limit, len(doc))):
            text += doc[i].get_text() + "\n"
        doc.close()
    except Exception as e:
        logger.error(f"Error reading PDF {path}: {e}")
    return text

def extract_text_pptx(path, limit=10):
    if not Presentation:
        return ""
    text = ""
    try:
        prs = Presentation(path)
        for i, slide in enumerate(prs.slides):
            if i >= limit: break
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        logger.error(f"Error reading PPTX {path}: {e}")
    return text

def score_folder(folder_name, audio_date):
    # folder examples: "0318-0319 강의자료(김승준)", "0320-0324"
    # Matches patterns like 0318-0319
    match = re.search(r'(\d{2})(\d{2})-(\d{2})(\d{2})', folder_name)
    if match:
        m1, d1, m2, d2 = map(int, match.groups())
        try:
            date1 = datetime.date(audio_date.year, m1, d1)
            date2 = datetime.date(audio_date.year, m2, d2)
            
            # If strictly within the bounds
            if date1 <= audio_date <= date2:
                return 1.0
            
            # Close proximity
            diff = min(abs((audio_date - date1).days), abs((audio_date - date2).days))
            if diff <= 7:
                return 0.6
                
            return 0.1
        except Exception:
            pass
            
    # Matches N월
    match_month = re.search(r'(\d+)월', folder_name)
    if match_month:
        if int(match_month.group(1)) == audio_date.month:
            return 0.4

    return 0.0

def find_candidates(sync_folder, audio_date):
    candidates = []
    if not os.path.exists(sync_folder):
        logger.warning(f"Sync folder not found: {sync_folder}")
        return candidates
        
    for root, dirs, files in os.walk(sync_folder):
        folder_name = os.path.basename(root)
        parent_name = os.path.basename(os.path.dirname(root))
        
        # Best guess from either current folder name or parent
        f_score = max(score_folder(folder_name, audio_date), score_folder(parent_name, audio_date))
        
        # Don't waste time on completely irrelevant folders if they score 0 and we are strict
        # But let's be generous for the MVP to avoid missing files.
        for f in files:
            ext = f.lower().split('.')[-1]
            if ext in ['pdf', 'pptx'] and not f.startswith('~'):
                candidates.append({
                    "path": os.path.join(root, f),
                    "filename": f,
                    "folder_score": f_score
                })
    return candidates

def retrieve_materials(sync_folder, audio_date_str, transcript):
    audio_date = parse_date(audio_date_str)
    
    logger.info(f"Scanning '{sync_folder}' for materials around {audio_date}...")
    candidates = find_candidates(sync_folder, audio_date)
    
    if not candidates:
        logger.info("No candidates found in sync folder.")
        return []

    # Sort and take top 10 to avoid heavy PDF parsing on hundreds of files
    candidates.sort(key=lambda x: x['folder_score'], reverse=True)
    top_candidates = candidates[:10]
    
    logger.info(f"Selected {len(top_candidates)} raw candidate files based on folder heuristics.")

    corpus = [transcript]
    
    for c in top_candidates:
        logger.info(f"Extracting text from: {c['filename']}")
        if c['path'].lower().endswith('.pdf'):
            c['text'] = extract_text_pdf(c['path'], limit=10)
        elif c['path'].lower().endswith('.pptx'):
            c['text'] = extract_text_pptx(c['path'], limit=10)
        else:
            c['text'] = ""
            
        corpus.append(c['filename'] + " \n " + c['text'])

    logger.info("Ranking candidates based on TF-IDF similarity to transcript...")
    vectorizer = TfidfVectorizer(max_features=2000, stop_words='english')
    try:
        tfidf_matrix = vectorizer.fit_transform(corpus)
        # Similarity of transcript (index 0) with all candidates (indices 1 to N)
        similarity = cosine_similarity(tfidf_matrix[0:1], tfidf_matrix[1:]).flatten()
    except ValueError:
        logger.warning("Corpus too small or empty for TF-IDF.")
        return []

    results = []
    for i, score in enumerate(similarity):
        c = top_candidates[i]
        # Weighted score: 30% rule-based folder score, 70% semantic content score
        final_score = (0.30 * c['folder_score']) + (0.70 * score)
        results.append({
            "path": c['path'],
            "filename": c['filename'],
            "score": final_score,
            "text": c['text'][:1000] # trim attached text payload for prompt
        })
        
    results.sort(key=lambda x: x['score'], reverse=True)
    
    # Thresholding & Output
    THRESHOLD = 0.05  # A bit tolerant for MVP
    best = [r for r in results if r['score'] >= THRESHOLD]
    top_3 = best[:3]
    
    if top_3:
        logger.info("Target materials selected:")
        for r in top_3:
            logger.info(f" - {r['filename']} (Score: {r['score']:.3f})")
    else:
        logger.info("No materials passed the relevance threshold.")
        
    return top_3
